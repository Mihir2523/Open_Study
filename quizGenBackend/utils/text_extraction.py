import fitz  # PyMuPDF for PDFs
from pptx import Presentation
from pathlib import Path


class FileTextExtractor:
    """
    Utility class for extracting text from different file types
    (PDF, PPT, PPTX).
    """

    SUPPORTED_EXTENSIONS = {".pdf", ".ppt", ".pptx"}

    @staticmethod
    def extract_text(file_path: str) -> str:
        """
        Extracts text from a file based on its extension.
        :param file_path: Path to the file
        :return: Extracted text as a string
        """
        ext = Path(file_path).suffix.lower()

        if ext not in FileTextExtractor.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {ext}")

        if ext == ".pdf":
            return FileTextExtractor._extract_pdf(file_path)
        elif ext in (".ppt", ".pptx"):
            return FileTextExtractor._extract_ppt(file_path)

    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        """Extract text from PDF using PyMuPDF."""
        text = []
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text.append(page.get_text())
        return "\n".join(text).strip()

    @staticmethod
    def _extract_ppt(file_path: str) -> str:
        """Extract text from PPT/PPTX using python-pptx."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text.append(para.text)
        return "\n".join(text).strip()
